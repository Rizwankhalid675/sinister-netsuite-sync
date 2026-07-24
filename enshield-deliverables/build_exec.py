#!/usr/bin/env python3
"""Build the 4-slide EXEC version of the Enshield System Blueprint."""
import importlib.util, sys
spec = importlib.util.spec_from_file_location("bd", "build_deck.py")
# We only want the helpers, not the module-level slide building. So inline instead:
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG_DARK=RGBColor(0x0E,0x14,0x1B); BG_PANEL=RGBColor(0x16,0x20,0x2B)
TEAL=RGBColor(0x16,0xB8,0xC4); TEAL_DK=RGBColor(0x0E,0x7C,0x86)
SILVER=RGBColor(0xC7,0xCE,0xD6); WHITE=RGBColor(0xF4,0xF7,0xFA)
MUTED=RGBColor(0x8A,0x97,0xA3); ACCENT=RGBColor(0x2E,0xD5,0xE0)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=BG_DARK; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def box(s,x,y,w,h,fill=None,line=None,line_w=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(line_w)
    shp.shadow.inherit=False; return shp
def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(6)
        if isinstance(para,tuple): para=[para]
        for (t,sz,col,bold) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=col
            r.font.bold=bold; r.font.name="Segoe UI"
    return tb
def kicker(s,txt):
    box(s,0.6,0.55,0.12,0.55,fill=TEAL); text(s,0.85,0.5,11,0.7,[[(txt,26,WHITE,True)]])
def bullets(s,x,y,w,h,items,sz=16):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(10)
        r=p.add_run(); r.text="▸  "; r.font.size=Pt(sz); r.font.color.rgb=TEAL; r.font.bold=True; r.font.name="Segoe UI"
        if isinstance(it,tuple):
            h_,rest=it
            r2=p.add_run(); r2.text=h_; r2.font.size=Pt(sz); r2.font.color.rgb=WHITE; r2.font.bold=True; r2.font.name="Segoe UI"
            r3=p.add_run(); r3.text=rest; r3.font.size=Pt(sz); r3.font.color.rgb=SILVER; r3.font.name="Segoe UI"
        else:
            r2=p.add_run(); r2.text=it; r2.font.size=Pt(sz); r2.font.color.rgb=SILVER; r2.font.name="Segoe UI"

# SLIDE 1 — TITLE
s=slide()
sh=s.shapes.add_shape(MSO_SHAPE.HEXAGON,Inches(5.66),Inches(1.1),Inches(2.0),Inches(2.2))
sh.fill.solid(); sh.fill.fore_color.rgb=BG_PANEL; sh.line.color.rgb=TEAL; sh.line.width=Pt(2.5); sh.shadow.inherit=False
text(s,5.66,1.7,2.0,1.0,[[("EN",40,TEAL,True)]],align=PP_ALIGN.CENTER)
text(s,1,3.7,11.333,1.2,[[("ENSHIELD REVAMP — EXEC BRIEF",40,WHITE,True)]],align=PP_ALIGN.CENTER)
text(s,1,4.75,11.333,0.7,[[("The 4-minute version — what it is, what's wrong, what we do",19,ACCENT,False)]],align=PP_ALIGN.CENTER)
box(s,4.6,5.75,4.13,0.02,fill=TEAL_DK)
text(s,1,5.95,11.333,0.6,[[("Prepared for Mike",15,MUTED,False)]],align=PP_ALIGN.CENTER)

# SLIDE 2 — WHAT IT IS (screenshot)
s=slide(); kicker(s,"What We Have Today")
box(s,1.81,1.29,9.72,5.27,fill=BG_PANEL,line=TEAL_DK,line_w=1.5)
s.shapes.add_picture("mock_dashboard.png",Inches(1.87),Inches(1.35),Inches(9.6),Inches(5.15))
text(s,0.85,6.65,11.6,0.5,[[("Insanelab-built dashboard pulling Shopify data via Gadget — Activity, KPI cards, client table.",14,MUTED,False)]])

# SLIDE 3 — WHAT'S WRONG + WHAT WE DO
s=slide(); kicker(s,"The Problem & The Plan")
box(s,0.85,1.5,5.7,4.9,fill=BG_PANEL)
text(s,1.15,1.75,5.1,0.5,[[("What's holding it back",18,ACCENT,True)]])
bullets(s,1.15,2.4,5.2,3.7,[
    ("Efficiency — ","heavy Shopify reads on every load; no caching layer"),
    ("Stale surfaces — ","KPIs recomputed client-side, slow to render"),
    ("Hard to extend — ","new features mean touching raw Gadget data each time"),
    ("Design drift — ","dashboard doesn't match the Enshield brand site"),
])
box(s,6.75,1.5,5.7,4.9,fill=BG_PANEL,line=TEAL_DK,line_w=1.2)
text(s,7.05,1.75,5.1,0.5,[[("What we do about it",18,TEAL,True)]])
bullets(s,7.05,2.4,5.2,3.7,[
    ("Aggregate in Gadget — ","pre-compute KPIs, dashboard just reads"),
    ("Cache + incremental sync — ","cut redundant Shopify calls"),
    ("Modular feature slots — ","add surfaces without schema surgery"),
    ("Re-skin to brand — ","match the Enshield public-site language"),
])

# SLIDE 4 — NEXT / ASK
s=slide(); kicker(s,"What I Need From You")
box(s,0.85,1.6,11.6,3.0,fill=BG_PANEL,line=TEAL_DK,line_w=1.2)
bullets(s,1.25,1.95,10.9,2.6,[
    ("Feature wishlist — ","the exact surfaces you want added to the revamped dashboard"),
    ("Access confirmed — ","Gadget + DigitalOcean + Shopify (done — full audit complete)"),
    ("Priority call — ","efficiency-first refactor, or new features first?"),
],sz=17)
text(s,0.85,5.0,11.6,0.6,[[("Phased plan: 1) aggregate & cache  →  2) re-skin  →  3) new features",17,ACCENT,True)]],align=PP_ALIGN.CENTER)
text(s,0.85,5.9,11.6,0.6,[[("Full 11-slide blueprint available on request.",14,MUTED,False)]],align=PP_ALIGN.CENTER)

out=sys.argv[1] if len(sys.argv)>1 else "Enshield-Exec-Brief.pptx"
prs.save(out); print("SAVED",out,"with",len(prs.slides._sldIdLst),"slides")
