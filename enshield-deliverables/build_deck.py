#!/usr/bin/env python3
"""Build the Enshield System Blueprint PowerPoint deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand palette (Enshield dark / teal / silver) ----
BG_DARK   = RGBColor(0x0E, 0x14, 0x1B)
BG_PANEL  = RGBColor(0x16, 0x20, 0x2B)
TEAL      = RGBColor(0x16, 0xB8, 0xC4)
TEAL_DK   = RGBColor(0x0E, 0x7C, 0x86)
SILVER    = RGBColor(0xC7, 0xCE, 0xD6)
WHITE     = RGBColor(0xF4, 0xF7, 0xFA)
MUTED     = RGBColor(0x8A, 0x97, 0xA3)
ACCENT    = RGBColor(0x2E, 0xD5, 0xE0)
GREEN     = RGBColor(0x3F, 0xB9, 0x50)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG_DARK; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs, str):
        runs = [[(runs, 18, WHITE, False)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        if isinstance(para, tuple): para = [para]
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Segoe UI"
    return tb

def bullets(s, x, y, w, h, items, sz=15, col=SILVER, gap=8):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "▸  "; r.font.size = Pt(sz); r.font.color.rgb = TEAL; r.font.bold = True
        if isinstance(it, tuple):
            head, rest = it
            r2 = p.add_run(); r2.text = head; r2.font.size = Pt(sz); r2.font.color.rgb = WHITE; r2.font.bold = True
            r3 = p.add_run(); r3.text = rest; r3.font.size = Pt(sz); r3.font.color.rgb = col
            for rr in (r2, r3): rr.font.name = "Segoe UI"
        else:
            r2 = p.add_run(); r2.text = it; r2.font.size = Pt(sz); r2.font.color.rgb = col; r2.font.name = "Segoe UI"
        r.font.name = "Segoe UI"
    return tb

def picture(s, path, x, y, w, h):
    # framed screenshot with teal border
    b = box(s, x-0.06, y-0.06, w+0.12, h+0.12, fill=BG_PANEL, line=TEAL_DK, line_w=1.5)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return b

def shot_slide(kick, path, caption):
    s = slide(); kicker(s, kick)
    # 1200x720 => 5:3; render at 9.6 x 5.76 centered
    picture(s, path, 1.87, 1.35, 9.6, 5.15)
    text(s, 0.85, 6.65, 11.6, 0.5, [[(caption, 14, MUTED, False)]])
    return s

def kicker(s, txt):
    box(s, 0.6, 0.55, 0.12, 0.55, fill=TEAL)
    text(s, 0.85, 0.5, 11, 0.7, [[(txt, 26, WHITE, True)]])

def footer(s, n):
    text(s, 0.6, 7.02, 8, 0.4, [[("ENSHIELD  ·  System Blueprint", 10, MUTED, False)]])
    text(s, 12.0, 7.02, 0.9, 0.4, [[(str(n), 10, MUTED, False)]], align=PP_ALIGN.RIGHT)

# ============================================================ SLIDE 1 — TITLE
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=BG_DARK)
# shield mark
sh = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(5.66), Inches(1.15), Inches(2.0), Inches(2.2))
sh.fill.solid(); sh.fill.fore_color.rgb = BG_PANEL; sh.line.color.rgb = TEAL; sh.line.width = Pt(2.5)
sh.shadow.inherit = False
text(s, 5.66, 1.75, 2.0, 1.0, [[("EN", 40, TEAL, True)]], align=PP_ALIGN.CENTER)
text(s, 1.0, 3.7, 11.333, 1.2, [[("ENSHIELD SYSTEM BLUEPRINT", 44, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, 1.0, 4.75, 11.333, 0.7,
     [[("How Gadget · DigitalOcean · Shopify · Dashboard fit together", 20, ACCENT, False)]],
     align=PP_ALIGN.CENTER)
box(s, 4.6, 5.75, 4.13, 0.02, fill=TEAL_DK)
text(s, 1.0, 5.95, 11.333, 0.6, [[("Findings & Architecture — prepared for Mike", 15, MUTED, False)]],
     align=PP_ALIGN.CENTER)

# ============================================================ SLIDE 2 — TL;DR
s = slide(); kicker(s, "The Big Picture (TL;DR)")
box(s, 0.85, 1.5, 11.6, 1.55, fill=BG_PANEL, line=TEAL_DK)
text(s, 1.15, 1.65, 11.0, 1.3,
     [[("We traced the whole Enshield stack end-to-end. There are ", 16, SILVER, False),
       ("two separate front-end surfaces", 16, TEAL, True),
       (", one ", 16, SILVER, False), ("Gadget backend", 16, TEAL, True),
       (" holding the data & logic, a ", 16, SILVER, False), ("DigitalOcean droplet", 16, TEAL, True),
       (" running a parallel dashboard, and ", 16, SILVER, False), ("Shopify", 16, TEAL, True),
       (" as the source of truth for stores, orders & checkout.", 16, SILVER, False)]],
     anchor=MSO_ANCHOR.MIDDLE)
cards = [("SHOPIFY", "Source of truth", "Stores, orders, carts,\ncheckout, protection product"),
         ("GADGET", "Backend brain", "Syncs data, holds models,\nauto-generates the API"),
         ("DIGITALOCEAN", "Hosting / legacy", "Droplet running a\nparallel dashboard"),
         ("DASHBOARD", "Two front ends", "Embedded settings page +\nstandalone admin UI")]
x = 0.85
for title, sub, body in cards:
    box(s, x, 3.35, 2.72, 3.1, fill=BG_PANEL, line=TEAL_DK)
    box(s, x, 3.35, 2.72, 0.14, fill=TEAL)
    text(s, x+0.2, 3.62, 2.4, 0.5, [[(title, 15, WHITE, True)]])
    text(s, x+0.2, 4.08, 2.4, 0.4, [[(sub, 12, ACCENT, True)]])
    text(s, x+0.2, 4.6, 2.4, 1.7, [[(body, 12.5, MUTED, False)]])
    x += 2.9
footer(s, 2)

# ============================================================ SLIDE 3 — FOUR SYSTEMS
s = slide(); kicker(s, "The Four Systems")
data = [
    ("1  SHOPIFY", "Source of truth",
     ["Merchants install the Enshield app on their stores",
      "Owns shops, orders, line items, carts, checkout",
      "Protection product injected at checkout",
      "Enshield reads this — it does not replace Shopify"]),
    ("2  GADGET", "Backend / brain",
     ["App: enshield-shipping-protection (built on Gadget)",
      "Continuously syncs Shopify data into its models",
      "Models: shopifyShop, shopifyOrder, shopifyCart,",
      "shippingInsuranceProduct / Setting, sync, session",
      "Auto-generates the access-controlled API"]),
    ("3  FRONT ENDS", "Two surfaces (key finding)",
     ["A) Embedded Shopify app — 1 Polaris settings page",
      "B) Standalone admin dashboard — the dark UI with",
      "Value in Transit, Clients, Claims, Reports",
      "B is the real revamp target"]),
    ("4  DIGITALOCEAN", "Hosting / legacy",
     ["Droplet runs a version of the dashboard/backend",
      "Accessed via SSH to confirm prod setup",
      "Revamp decision: keep on DO, move to Gadget,",
      "or hybrid hosting"]),
]
positions = [(0.85, 1.45), (6.95, 1.45), (0.85, 4.35), (6.95, 4.35)]
for (title, sub, items), (px, py) in zip(data, positions):
    box(s, px, py, 5.5, 2.75, fill=BG_PANEL, line=TEAL_DK)
    text(s, px+0.25, py+0.15, 5.0, 0.4, [[(title, 16, TEAL, True)]])
    text(s, px+0.25, py+0.55, 5.0, 0.35, [[(sub, 11.5, ACCENT, True)]])
    bullets(s, px+0.25, py+0.95, 5.0, 1.7, items, sz=11, gap=3)
footer(s, 3)

# ============================================================ SLIDE 4 — DATA FLOW
s = slide(); kicker(s, "How Data Flows")
def node(x, y, w, h, title, sub, fill=BG_PANEL, edge=TEAL_DK, tcol=WHITE):
    box(s, x, y, w, h, fill=fill, line=edge, line_w=1.5)
    text(s, x, y+0.14, w, 0.45, [[(title, 15, tcol, True)]], align=PP_ALIGN.CENTER)
    if sub:
        text(s, x, y+0.62, w, 0.7, [[(sub, 10.5, MUTED, False)]], align=PP_ALIGN.CENTER)
def arrow(x, y, w, h, rot=0):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = TEAL; a.line.fill.background(); a.rotation = rot
    a.shadow.inherit = False

node(4.9, 1.35, 3.5, 0.95, "Merchant's Shopify Store", "installs app · checkout adds protection")
arrow(6.5, 2.4, 0.3, 0.5)
node(4.9, 3.0, 3.5, 1.05, "SHOPIFY", "orders · carts · shops · checkout", fill=BG_PANEL, edge=TEAL, tcol=TEAL)
arrow(6.5, 4.15, 0.3, 0.5)
text(s, 7.0, 4.18, 3.5, 0.4, [[("webhooks / sync  →", 11, ACCENT, True)]])
node(4.9, 4.65, 3.5, 1.15, "GADGET BACKEND",
     "models + auto-generated, access-controlled API", fill=BG_PANEL, edge=TEAL, tcol=TEAL)
# split to two front ends
arrow(2.75, 6.0, 0.3, 0.45, rot=315)
arrow(10.25, 6.0, 0.3, 0.45, rot=45)
node(0.75, 6.35, 4.5, 0.95, "A) Embedded Settings Page", "Polaris · inside Shopify admin · 1 screen")
node(8.1, 6.35, 4.5, 0.95, "B) Standalone Admin Dashboard", "Value in Transit · Clients · Claims · Reports")
# DO tag
box(s, 8.1, 5.55, 4.5, 0.55, fill=None, line=MUTED)
text(s, 8.1, 5.6, 4.5, 0.45, [[("hosted on / near DigitalOcean", 11, MUTED, True)]], align=PP_ALIGN.CENTER)
footer(s, 4)

# ============================================================ SLIDE 5 — KPI MAPPING
s = slide(); kicker(s, "Dashboard KPIs → Gadget Data")
rows = [("Clients", "count of connected shopifyShop records", "5"),
        ("Value in Transit", "Σ active shippingInsuranceProduct across open orders", "$1,178,054"),
        ("Open Claims", "claims count from the claims workflow", "0"),
        ("Latest Clients", "shopifyShop list w/ per-shop value, claims, status", "table")]
text(s, 0.85, 1.4, 11.6, 0.5,
     [[("Every number on the dark dashboard maps directly back to a Gadget model:", 15, SILVER, False)]])
y = 2.1
box(s, 0.85, y, 11.6, 0.55, fill=TEAL_DK)
text(s, 1.05, y+0.08, 3.0, 0.4, [[("KPI", 14, WHITE, True)]])
text(s, 4.2, y+0.08, 6.0, 0.4, [[("Comes from", 14, WHITE, True)]])
text(s, 10.6, y+0.08, 1.7, 0.4, [[("Example", 14, WHITE, True)]])
y += 0.55
for i, (k, src, ex) in enumerate(rows):
    fill = BG_PANEL if i % 2 == 0 else BG_DARK
    box(s, 0.85, y, 11.6, 0.82, fill=fill, line=TEAL_DK, line_w=0.5)
    text(s, 1.05, y+0.18, 3.1, 0.5, [[(k, 14, TEAL, True)]])
    text(s, 4.2, y+0.18, 6.2, 0.5, [[(src, 12.5, SILVER, False)]])
    text(s, 10.6, y+0.18, 1.7, 0.5, [[(ex, 13, GREEN, True)]])
    y += 0.82
footer(s, 5)

# ============================================================ SLIDE 6 — WHAT IT MEANS
s = slide(); kicker(s, "What This Means for the Revamp")
items = [
    ("The embedded Shopify app is tiny — ", "basically one settings screen. Easy to modernize."),
    ("The real revamp target is the standalone admin dashboard ", "(the dark UI). Rebuild it against the Gadget API so it always reflects live synced data."),
    ("Gadget is the backend we build on — ", "add aggregation/computed actions for the KPIs and expose clean, access-controlled reads."),
    ("DigitalOcean is a hosting decision — ", "keep, migrate to Gadget hosting, or hybrid."),
]
box(s, 0.85, 1.5, 11.6, 3.7, fill=BG_PANEL, line=TEAL_DK)
bullets(s, 1.25, 1.85, 10.8, 3.2, items, sz=17, gap=16)
footer(s, 6)

# ============================================================ SLIDE 7 — PLAN + OPEN Qs
s = slide(); kicker(s, "Phased Plan & Open Questions")
box(s, 0.85, 1.45, 6.7, 5.1, fill=BG_PANEL, line=TEAL_DK)
text(s, 1.1, 1.6, 6.2, 0.5, [[("Phased Plan", 18, TEAL, True)]])
phases = [
    ("Phase 1 — Architecture lock-in", "confirm target surface + hosting; reuse audit report as baseline"),
    ("Phase 2 — Backend (Gadget)", "aggregation actions for KPIs; clean access-controlled read API"),
    ("Phase 3 — Frontend revamp", "rebuild dashboard on Gadget API; modernize settings page"),
    ("Phase 4 — New features", "implement Brain's feature list"),
]
yy = 2.25
for h, b in phases:
    box(s, 1.1, yy, 0.1, 0.75, fill=TEAL)
    text(s, 1.35, yy, 6.0, 0.4, [[(h, 14, WHITE, True)]])
    text(s, 1.35, yy+0.36, 6.0, 0.5, [[(b, 12, MUTED, False)]])
    yy += 1.0

box(s, 7.85, 1.45, 4.6, 5.1, fill=BG_PANEL, line=TEAL)
text(s, 8.1, 1.6, 4.1, 0.5, [[("Need from Brain", 18, ACCENT, True)]])
qs = ["Feature list — what to add/change?",
      "Target surface — dashboard, embedded app, or both?",
      "Hosting — DigitalOcean, Gadget, or hybrid?",
      "Where the revamped code lives?"]
bullets(s, 8.1, 2.3, 4.1, 4.0, qs, sz=14, gap=18)
footer(s, 7)

# ===================================================== SCREENSHOT SLIDES
shot_slide("The Dashboard Today", "mock_dashboard.png",
           "Current Insanelab dashboard — Activity chart, KPI cards (Clients · Value in Transit · Open Claims), Latest Clients table. This is what we're revamping.")
shot_slide("The Public Site", "mock_homepage.png",
           "Enshield marketing homepage — the brand/design language the revamped dashboard should echo.")
shot_slide("The Gadget Backend", "mock_gadget.png",
           "Gadget apps list — per-merchant Shopify apps + the enshield-shipping-protection core. This is the real data source behind the dashboard.")

# ============================================================ SLIDE 8 — CLOSE
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=BG_DARK)
box(s, 4.6, 3.55, 4.13, 0.03, fill=TEAL)
text(s, 1.0, 2.7, 11.333, 0.9, [[("One system. Four moving parts.", 34, WHITE, True)]], align=PP_ALIGN.CENTER)
text(s, 1.0, 3.75, 11.333, 0.7,
     [[("Shopify feeds Gadget · Gadget powers the dashboards · DO hosts them.", 17, ACCENT, False)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 4.6, 11.333, 0.5, [[("Answer the four questions and we start building.", 15, MUTED, False)]],
     align=PP_ALIGN.CENTER)

import sys
_out = sys.argv[1] if len(sys.argv) > 1 else "Enshield-System-Blueprint.pptx"
prs.save(_out)
print("SAVED", _out, "with", len(prs.slides._sldIdLst), "slides")
